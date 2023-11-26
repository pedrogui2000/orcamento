#MODIFICAÇÂO PDF ORCAMENTO MEC

#Modulos importados
import tkinter as tk
from tkinter import Tk, ttk ,Label, Entry, Button, messagebox, Listbox, Text, filedialog, Toplevel, StringVar, END
from PIL import Image, ImageTk
from pptx import Presentation
import os
import openpyxl
from pptx.util import Inches
import locale
from datetime import datetime

#Classe principal da Interface
class InterfaceGrafica:
    #Gerando Janela principal da Interface
    def __init__(self, janela):
        self.janela = janela
        self.janela.title("Orçamento MEC")
        self.num_slides = 0  # Inicializa o número de slides

        # Estilo das caixas
        style = ttk.Style(self.janela)
        style.theme_use("vista")

        # Nome Empresa
        self.rotulo_nome_empresa = Label(janela, text="Nome da Empresa Contratante:")
        self.rotulo_nome_empresa.grid(row=0, column=0, pady=10, padx=10, sticky="e")

        self.caixa_texto_nome_empresa = ttk.Entry(janela, style="TEntry")
        self.caixa_texto_nome_empresa.grid(row=0, column=1, pady=10, padx=10, sticky="w")

        # Nome do Serviço
        self.rotulo_nome_servico = Label(janela, text="Nome do Serviço:")
        self.rotulo_nome_servico.grid(row=1, column=0, pady=10, padx=10, sticky="e")

        self.caixa_texto_nome_servico = ttk.Entry(janela, style="TEntry")
        self.caixa_texto_nome_servico.grid(row=1, column=1, pady=10, padx=10, sticky="w")

        # Id do Serviço

        self.rotulo_id_contrato = Label(janela, text="Digite o ID do Orçamento:")
        self.rotulo_id_contrato.grid(row=2, column=0, pady=10, padx=10, sticky="e")

        self.caixa_id_contrato = ttk.Entry(janela, style="TEntry")
        self.caixa_id_contrato.grid(row=2, column=1, pady=10, padx=10, sticky="w")

        # Prazo do Serviço

        self.rotulo_prazo_trampo = Label(janela, text="Prazo de execução do serviço:")
        self.rotulo_prazo_trampo.grid(row=3, column=0, pady=10, padx=10, sticky="e")

        self.caixa_prazo_trampo = ttk.Entry(janela, style="TEntry")
        self.caixa_prazo_trampo.grid(row=3, column=1, pady=10, padx=10, sticky="w")

        # Validade da Proposta

        self.rotulo_prazo_proposta = Label(janela, text="Validade da Proposta:")
        self.rotulo_prazo_proposta.grid(row=4, column=0, pady=10, padx=10, sticky="e")

        self.caixa_prazo_proposta = ttk.Entry(janela, style="TEntry")
        self.caixa_prazo_proposta.grid(row=4, column=1, pady=10, padx=10, sticky="w")

        # Elementos para a lista de itens

        self.lista_itens = []

        #Caixa onde os itens são exibidos

        self.lista_box_itens = Listbox(janela, height=5, selectmode="SINGLE", width=50)
        self.lista_box_itens.grid(row=5, column=0, columnspan=2, pady=10, padx=10)

        #Botões de adicionar e remover itens da lista

        self.botao_remover_item = Button(janela, text="Remover Item", command=self.remover_item)
        self.botao_remover_item.grid(row=6, column=0, pady=5, padx=10, sticky="e")

        self.botao_adicionar_item = Button(janela, text="Adicionar Item", command=self.adicionar_item)
        self.botao_adicionar_item.grid(row=6, column=1, pady=10, padx=10, sticky="w")

        # Adiciona desconto
        self.rotulo_desconto = Label(janela, text="Deseja adicionar desconto?")
        self.rotulo_desconto.grid(row=7, column=0, pady=10, padx=10, sticky="e")

        self.caixa_desconto = ttk.Entry(janela, style="TEntry")
        self.caixa_desconto.grid(row=7, column=1, pady=10, padx=10, sticky="w")

        # Obs 1

        self.rotulo_obs1 = tk.Label(janela, text="Observação 1:")
        self.rotulo_obs1.grid(row=8, column=0, pady=10, padx=10, sticky="e")

        self.caixa_obs1 = tk.Text(janela, height=10, width=80, wrap="word")
        self.caixa_obs1.grid(row=8, column=1, pady=5, padx=5, sticky="w")

        # Obs 2

        self.rotulo_obs2 = tk.Label(janela, text="Observação 2:")
        self.rotulo_obs2.grid(row=9, column=0, pady=10, padx=10, sticky="e")

        self.caixa_obs2 = tk.Text(janela, height=5, width=80, wrap="word")
        self.caixa_obs2.grid(row=9, column=1, pady=5, padx=5, sticky="w")
        
        # Botao adicionar imagem

        self.botao_selecionar_imagens= Button(janela, text="Selecionar imagens", command=self.selecionar_imagens)
        self.botao_selecionar_imagens.grid(row=4, column=2, columnspan=2, pady=20, padx=10 , )

       # Etiqueta para mostrar o nome das imagens selecionadas
        self.etiqueta_imagens_selecionadas = Label(janela, text="Imagens Selecionadas: ")
        self.etiqueta_imagens_selecionadas.grid(row=5, column=2, columnspan=2, pady=10, padx=10)


        # Botão para gerar orçamento
        self.botao_gerar_orcamento = Button(janela, text="Gerar Orçamento", command=self.gerar_orcamento)
        self.botao_gerar_orcamento.grid(row=6, column=2, columnspan=2, pady=20, padx=10 , )

         # Adicione o atributo arquivos_selecionados
        self.arquivos_selecionados = []

    def adicionar_item(self):
        # Abrir uma nova janela para adicionar um item
        janela_item = Toplevel(self.janela)
        janela_item.title("Adicionar Item")

        # Rótulos e Caixas de Texto para cada campo do item
        rotulo_item = Label(janela_item, text="Item:")
        rotulo_item.grid(row=0, column=0, padx=5, pady=5)

        caixa_texto_item = ttk.Entry(janela_item, style="TEntry", width=100)
        caixa_texto_item.grid(row=0, column=1, padx=5, pady=5)

        rotulo_descricao = Label(janela_item, text="Descrição:")
        rotulo_descricao.grid(row=1, column=0, padx=5, pady=5)

        caixa_texto_descricao = ttk.Entry(janela_item, style="TEntry",width=100)
        caixa_texto_descricao.grid(row=1, column=1, padx=5, pady=5)

        rotulo_valor_unit = Label(janela_item, text="Valor Unit:")
        rotulo_valor_unit.grid(row=2, column=0, padx=5, pady=5)

        caixa_texto_valor_unit = ttk.Entry(janela_item, style="TEntry", width=100)
        caixa_texto_valor_unit.grid(row=2, column=1, padx=5, pady=5)

        rotulo_qtd = Label(janela_item, text="Qtde:")
        rotulo_qtd.grid(row=3, column=0, padx=5, pady=5)

        caixa_texto_qtd = ttk.Entry(janela_item, style="TEntry",width=100)
        caixa_texto_qtd.grid(row=3, column=1, padx=0, pady=5)

        # Configurar a largura da nova janela
        janela_item.geometry("800x200")  # <--- MUDAR CONFORME NECESSARIO AQUI
        
        # Função para adicionar o item à lista
        def adicionar_item_lista():
            descricao = caixa_texto_descricao.get()
            valor_unit = caixa_texto_valor_unit.get()
            qtd = caixa_texto_qtd.get()

            # Verificar se todos os campos estão preenchidos
            if descricao and valor_unit and qtd:
                item_str = f"{descricao} - {valor_unit} - {qtd}"
                self.lista_itens.append((descricao, valor_unit, qtd))
                self.lista_box_itens.insert(END, item_str)
                janela_item.destroy()
            else:
                messagebox.showerror("Erro", "Preencha todos os campos.")

        # Botão para adicionar o item à lista
        botao_adicionar = Button(janela_item, text="Adicionar", command=adicionar_item_lista)
        botao_adicionar.grid(row=5, columnspan=2, pady=10)

    def remover_item(self):
        # Remover item selecionado da lista
        selecionado = self.lista_box_itens.curselection()
        if selecionado:
            index = int(selecionado[0])
            del self.lista_itens[index]
            self.lista_box_itens.delete(selecionado)

    def selecionar_imagens(self):
        #Função para buscar as imagens no computador do usuario para adicionar ao slide
        filenames = filedialog.askopenfilenames(
            initialdir=os.getcwd(),
            title='Selecione as imagens',
            filetypes=(('Arquivos de imagem', '*.jpg *.png'), ('Todos os arquivos', '*.*'))
        )

        # Verificar se pelo menos uma imagem foi selecionada
        if not filenames:
            messagebox.showerror("Erro", "Selecione ao menos uma imagem.")
            return

        # Limatador de caractere exibido na interface ( Nome das Imagens selecionadas)
        nomes_imagens = [f"- {os.path.basename(filename)[:50] + '...' if len(os.path.basename(filename)) > 50 else os.path.basename(filename)}" for filename in filenames]

        # Adicionar os nomes das imagens à etiqueta
        if hasattr(self, 'nomes_imagens_selecionadas'):
            self.nomes_imagens_selecionadas.extend(nomes_imagens)
        else:
            self.nomes_imagens_selecionadas = nomes_imagens

        # Atualizar a etiqueta com todos os nomes das imagens
        nomes_imagens_str = "\n".join(self.nomes_imagens_selecionadas)
        self.etiqueta_imagens_selecionadas.config(text=f"Imagens Selecionadas:\n{nomes_imagens_str}")

        # Adicione a seguinte linha para garantir que self.arquivos_selecionados contenha os nomes dos arquivos
        self.arquivos_selecionados = filenames

    def adicionar_ao_slide(self, caminho_pptx_modificado):
        # Adicione o conteúdo do método adicionar_ao_slide
        presentation = Presentation(caminho_pptx_modificado)

        slide_index = 2  # Slide onde as imagens vão ser adicionadas (subtrai 1 porque os índices começam em 0)
        slide = presentation.slides[slide_index]

        left = top = Inches(2)  # Ajuste a posição conforme necessário

        for arquivo in self.arquivos_selecionados:
            picture = slide.shapes.add_picture(arquivo, left, top, width=Inches(2), height=Inches(2))
            left += picture.width + Inches(0.5)  # Espaçamento entre as imagens

        # Salvar as alterações no arquivo pptx
        presentation.save(caminho_pptx_modificado)

    def gerar_orcamento(self):
        #Função principal de adicionar as alterações no arquivo pptx (orcamento pptx novo)
        nome_empresa = self.caixa_texto_nome_empresa.get().strip()
        nome_servico = self.caixa_texto_nome_servico.get().strip()
        id_contrato = self.caixa_id_contrato.get().strip()
        desconto = self.caixa_desconto.get().strip()
        prazotrampo = self.caixa_prazo_trampo.get().strip()
        obs1 = self.caixa_obs1.get("1.0", tk.END).strip()
        obs2 = self.caixa_obs2.get("1.0", tk.END).strip()
        valeday = self.caixa_prazo_proposta.get().strip()

        
        if nome_empresa and nome_servico:
            dados_empresa = self.ler_planilha_excel(nome_empresa)

            if dados_empresa:
                caminho_pptx_modificado = f"{nome_empresa} ORÇAMENTO.pptx"
                self.preencher_apresentacao_powerpoint(
                    nome_empresa, nome_servico, dados_empresa, caminho_pptx_modificado, id_contrato, desconto, prazotrampo, obs1, obs2, valeday)
                
                # Adicione a seguinte linha para adicionar as imagens ao slide
                self.adicionar_ao_slide(caminho_pptx_modificado)
                
                messagebox.showinfo("Sucesso", "Orçamento gerado com sucesso!")
            else:
                messagebox.showerror("Erro", "Empresa não encontrada na planilha.")
        else:
            messagebox.showerror("Erro", "Digite o nome da empresa contratante e o nome do serviço.")
            
    def ler_planilha_excel(self, nome_empresa):
        try:
            #Localiza o arquivo xlsm dentro da pasta
            caminho_excel = os.path.join(os.path.dirname(__file__), "AddClient.xlsm")

            # Ler o arquivo xlsm usando a biblioteca openpyxl
            wb = openpyxl.load_workbook(caminho_excel, read_only=True)
            planilha = wb['Clientes']

            #Fa a busca dentro da planilha atras dos dados 
            for row in range(1, planilha.max_row + 1):
                if planilha.cell(row=row, column=1).value and planilha.cell(row=row, column=1).value.lower() == nome_empresa.lower():
                    cabecalhos = ["Nome da Empresa", "Cnpj", "enderecocont", "Telefone", "Email", "Contato"]
                    dados_empresa = dict(zip(cabecalhos, [planilha.cell(row=row, column=col).value for col in range(1, planilha.max_column + 1)]))
                    return dados_empresa

            print("Empresa não encontrada na planilha.")
            return None
        except Exception as e:
            print(f"Erro ao ler planilha: {e}")
            return None
    
    def preencher_apresentacao_powerpoint(self, nome_empresa, nome_servico, dados_empresa, caminho_pptx_novo, id_contrato , desconto, prazotrampo , obs1, obs2, valeday):
        #Adicione as alterações no arquivo pptxx
        try:
            os.chdir(os.path.dirname(__file__))

            caminho_pptx = "orcamento.pptx"

            if os.path.exists(caminho_pptx):
                presentation = Presentation(caminho_pptx)

                for i, slide in enumerate(presentation.slides):
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            self.substituir_palavras_chave(shape.text_frame, nome_empresa, nome_servico, dados_empresa, id_contrato, desconto, prazotrampo, obs1, obs2, valeday)
                                              
                    # Adicionar data a partir da terceira página
                    if i >= 2:
                        self.adicionar_data_atualizacao(slide)

                # Salvar o PowerPoint modificado
                presentation.save(caminho_pptx_novo)

            else:
                print("Arquivo PowerPoint original não encontrado.")
        except Exception as e:
            print(f"Erro ao preencher apresentação PowerPoint: {e}")

    def substituir_palavras_chave(self, text_frame, nome_empresa, nome_servico, dados_empresa,id_contrato, desconto, prazotrampo, obs1, obs2, valeday):
        #Indendifica as palavras chaves dentro da planilha e adicione as novas
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                texto_original = run.text

                # Substituir individualmente cada item
                texto_original = texto_original.replace("NOMECONT", nome_empresa)
                texto_original = texto_original.replace("cnpjcont", dados_empresa["Cnpj"])
                texto_original = texto_original.replace("enderecocont", dados_empresa["enderecocont"])
                texto_original = texto_original.replace("telcont", dados_empresa["Telefone"])
                texto_original = texto_original.replace("emailcont", dados_empresa["Email"])
                texto_original = texto_original.replace("contatocont", dados_empresa["Contato"])
                texto_original = texto_original.replace("nametrampo", nome_servico)
                texto_original = texto_original.replace("idcontrato1", "N°"+ id_contrato) 
                texto_original = texto_original.replace("desconto","- "+"%"+(desconto)) 
                texto_original = texto_original.replace("daysexec", (prazotrampo) + " Dias") 
                texto_original = texto_original.replace("obs1", (obs1))
                texto_original = texto_original.replace("obs2", (obs2))
                texto_original = texto_original.replace("valeday", (valeday) + " Dias")
        
                run.text = texto_original

    def adicionar_data_atualizacao(self, slide):
        #Adiciona a data altomaticamente no slide substituindo a palabra chave "dateday"
        for shape in slide.shapes:
            if shape.has_text_frame:
                self.substituir_data_atualizacao(shape.text_frame)

    def substituir_data_atualizacao(self, text_frame):
        # Muda o idioma do mes para português do Brasil (pt_BR)
        locale.setlocale(locale.LC_TIME, 'pt_BR')

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if "dateday" in run.text:
                    data_atual = datetime.now().strftime("%d de %B de %Y")
                    
                    # Formatação do nome do mês para o português
                    data_atual = data_atual.capitalize()  # Primeira letra maiúscula
                    data_atual = data_atual.replace('De', 'de')  # Pequena correção

                    run.text = run.text.replace("dateday", data_atual)

    def inserir_texto(self):
        texto = self.entry_text.get()
if __name__ == "__main__":
    janela_principal = tk.Tk()
    app = InterfaceGrafica(janela_principal)
    janela_principal.mainloop()